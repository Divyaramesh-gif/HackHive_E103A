import os
import PyPDF2
import docx2txt
import numpy as np
import faiss
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure Gemini
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    print("Warning: GOOGLE_API_KEY not found. Please set it in .env")
else:
    genai.configure(api_key=api_key)

def extract_text_from_pdf(file_path):
    text = ""
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
    return text

def extract_text_from_docx(file_path):
    try:
        return docx2txt.process(file_path)
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
    return text

def extract_text(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
        
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file type")

def chunk_text(text, chunk_size=500):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunks.append(" ".join(words[i:i+chunk_size]))
    return chunks

def build_vector_store(chunks):
    try:
        # Use Gemini for embeddings
        result = genai.embed_content(
            model="models/text-embedding-004",
            content=chunks,
            task_type="retrieval_document",
            title="Custom Document"
        )
        embeddings = np.array(result['embedding'])
        
        dimension = embeddings.shape[1]
        index = faiss.IndexFlatL2(dimension)
        index.add(embeddings)
        return index, chunks
    except Exception as e:
        raise RuntimeError(f"Failed to generate embeddings: {e}")

def retrieve(query, index, chunks, top_k=3):
    try:
        # Embed query
        result = genai.embed_content(
            model="models/text-embedding-004",
            content=query,
            task_type="retrieval_query"
        )
        query_vec = np.array([result['embedding']])
        
        distances, indices = index.search(query_vec, top_k)
        return [chunks[i] for i in indices[0]]
    except Exception as e:
        print(f"Retrieval failed: {e}")
        return []

def answer_question(query, index, chunks):
    retrieved_chunks = retrieve(query, index, chunks, top_k=3)
    if not retrieved_chunks:
        return "Could not retrieve context to answer the question."
        
    context = " ".join(retrieved_chunks)
    
    prompt = f"""
    Answer the question based only on the following content. Do not invent information. Be accurate and unbiased.

    Content: 
    {context}

    Question: {query}
    """
    
    # Retry logic for generation
    import time
    max_retries = 3
    base_delay = 2
    
    for attempt in range(max_retries):
        try:
            model = genai.GenerativeModel('gemini-2.0-flash')
            response = model.generate_content(prompt)
            
            # Check if value is returned
            if not response.text:
                return "Error: Empty response from AI."
                
            return response.text
            
        except Exception as e:
            if "429" in str(e) or "quota" in str(e).lower() or "resource" in str(e).lower():
                if attempt < max_retries - 1:
                    wait_time = base_delay * (2 ** attempt)
                    print(f"Rate limit hit. Retrying in {wait_time}s...")
                    time.sleep(wait_time)
                    continue
                else:
                    return f"Error: API error: 429 (Rate Limit Exceeded). Please try again later."
            return f"Error generating answer: {e}"

def main():
    file_path = "example.pdf"
    
    if not os.path.exists(file_path):
        print(f"Sample file '{file_path}' not found. Please provide a valid file.")
        return

    print(f"Extracting text from {file_path}...")
    text = extract_text(file_path)
    if not text:
        print("No text extracted.")
        return

    print("Chunking text...")
    chunks = chunk_text(text, chunk_size=300)
    if not chunks:
        print("No chunks created.")
        return

    print("Building vector store (using Gemini Embeddings)...")
    try:
        index, chunks = build_vector_store(chunks)
    except Exception as e:
        print(f"Failed to build vector store: {e}")
        return

    query = "Explain the main process described in this document."
    print(f"Querying: {query}")
    answer = answer_question(query, index, chunks)
    print("Answer:", answer)

if __name__ == "__main__":
    main()